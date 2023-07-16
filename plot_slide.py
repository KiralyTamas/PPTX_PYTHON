from pptx.util import Inches, Pt
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
import csv

'''Grafikon diagramm'''

def plot(item,slide,pres):
    slide.shapes._spTree.remove(slide.shapes[1]._element)
    slide_width=pres.slide_width
    title=slide.shapes.title
    title.text = item['title']
    title.width = Inches(4)
    title.height = Inches(1)
    title.left=int((int(slide_width)-int(title.width))/2)
    title.top = Inches(0.5)
    csv_data=[]
    with open('sample.dat','r') as f:
        data=csv.reader(f)
        for i in data:
            csv_data.append(float(i[0].split(";")[0]))
            csv_data.append(float(i[0].split(";")[1]))
    csv_data=tuple(csv_data)
    chart_data=ChartData()
    chart_data.categories=[" "]
    chart_data.add_series("Diagramm",csv_data)
    x,y,cx,cy=Inches(2),Inches(2),Inches(6),Inches(4)
    chart=slide.shapes.add_chart(XL_CHART_TYPE.LINE,x,y,cx,cy,chart_data).chart
    category_axis = chart.category_axis
    category_axis_title = category_axis.axis_title
    category_axis_title.text_frame.text = item['configuration']['x-label']
    value_axis = chart.value_axis
    value_axis_title = value_axis.axis_title
    value_axis_title.text_frame.text = item['configuration']['y-label']
    return(slide)