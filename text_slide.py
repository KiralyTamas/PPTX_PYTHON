from pptx.util import Inches, Pt

'''Hosszabb szöveg diagramm'''

def text(item,slide,pres):
    slide_width=pres.slide_width
    title=slide.shapes.title
    title.text = item['title']
    title.width = Inches(4)
    title.height = Inches(1)
    title.left=int((int(slide_width)-int(title.width))/2)
    title.top = Inches(0.5)
    left = Inches(1.5)
    top = Inches(0.8)
    width = Inches(8)
    height = Inches(4)
    long_text=slide.shapes.add_textbox(left,top,width,height)
    text_frame=long_text.text_frame
    text_frame.word_wrap=True
    p=text_frame.add_paragraph()
    p.text=item['content']
    p.font.size=Pt(18)
    slide.shapes._spTree.remove(slide.shapes[1]._element)
    return(slide)