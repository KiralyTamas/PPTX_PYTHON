from pptx import Presentation
from pptx.util import Inches

'''Kép diagramm'''

def picture(item,slide,pres):
    slide_width=pres.slide_width
    title=slide.shapes.title
    title.text = item['title']
    title.width = Inches(4)
    title.height = Inches(1)
    title.left=int((int(slide_width)-int(title.width))/2)
    title.top = Inches(0.5)
    
    img_path="picture.png"
    img_width=Inches(6)
    img_height=Inches(4)
    left=int((int(slide_width)-int(img_width))/2)
    top=title.top+title.height
    slide.shapes.add_picture(img_path,left,top,img_width,img_height)
    slide.shapes._spTree.remove(slide.shapes[1]._element)
    return(slide)