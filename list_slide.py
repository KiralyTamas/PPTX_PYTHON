from pptx import Presentation
from pptx.util import Inches

'''Lista diagramm'''

def list(item,slide,pres):
    slide.shapes._spTree.remove(slide.shapes[1]._element)
    content=item['content']
    slide_width=pres.slide_width
    title=slide.shapes.title
    title.text = item['title']
    title.left=Inches(0.5)
    title.top = Inches(0.4)
    title.width = slide_width-(title.left*2)
    title.height = Inches(1.2)
    
    left=Inches(0.5)
    top=title.top+title.height
    width=slide_width-(left*2)
    height=Inches(4)
    list_box=slide.shapes.add_textbox(left,top,width,height)
    text_frame=list_box.text_frame
    for content_item in content:
        p=text_frame.add_paragraph()
        p.text=content_item['text']
        p.level=content_item['level']
    return(slide)