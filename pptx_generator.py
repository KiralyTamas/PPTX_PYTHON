from pptx import Presentation
from title_slide import *
from text_slide import *
from list_slide import *
from picture_slide import *
from plot_slide import *
import json

'''Itt kezdődik a pptx generálás.
    A json fájlból for ciklussal kikeresi a "type" adatokat
    és a megfelelő függvénybe adja be az adatokat
'''

def pptx_generator():
    pres=Presentation()
    slide_layout=pres.slide_layouts[0]
    with open('sample.json','r') as f:
        json_data=json.load(f)
    for item in json_data['presentation']:
        slide=pres.slides.add_slide(slide_layout)
        return_slide=eval(item['type'])(item=item,slide=slide,pres=pres)
        slide=return_slide
    pres.save('output.pptx')

pptx_generator()